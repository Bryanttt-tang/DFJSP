"""
create_fjsp_gantt_pptx.py

Creates a single PowerPoint slide with a Gantt-chart-like drawing for:
- 3 machines (M1, M2, M3)
- 5 jobs (J1..J5), each with 3 ordered operations (op1->op2->op3)
Each job has a single consistent color across its operations.
Each operation is drawn as an independent rectangle (shape) so you can add animations in PowerPoint.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# ---------- Parameters ----------
NUM_MACHINES = 3
MACHINES = ["M1", "M2", "M3"]
NUM_JOBS = 5
JOBS = [f"J{j+1}" for j in range(NUM_JOBS)]
OPS_PER_JOB = 3

# Visual layout (in inches)
SLIDE_WIDTH_IN = 13.3333
LEFT_MARGIN = 1.0
TOP_MARGIN = 1.0
ROW_HEIGHT = 0.8
ROW_GAP = 0.2
TIME_SCALE = 0.5  # inches per time unit (duration unit)
TIME_AXIS_OFFSET = LEFT_MARGIN + 1.2  # where time 0 starts (leaves space for machine labels)
TITLE_HEIGHT = 0.6

# A simple deterministic set of durations for reproducibility (or randomize as you like)
# durations[job_index][op_index] = duration value (time units)
durations = [
    [3, 2, 4],  # J1
    [2, 3, 2],  # J2
    [4, 1, 3],  # J3
    [1, 3, 2],  # J4
    [2, 2, 3],  # J5
]

# Machine required for each operation index (round-robin mapping for the example)
# op 0 -> machine 0, op 1 -> machine 1, op 2 -> machine 2
def required_machine_for_op(op_idx):
    return op_idx % NUM_MACHINES

# Job color map (distinct colors for each job)
JOB_COLORS = [
    (79, 129, 189),   # blue
    (192, 80, 77),    # red
    (155, 187, 89),   # green
    (128, 100, 162),  # purple
    (242, 169, 0)     # orange
]

# ---------- Simple earliest-possible scheduler ----------
# We'll schedule each job in job order; each job's operations must respect previous op completion.
# Each machine has its own available_time (initially 0).
def schedule_jobs():
    machine_available = [0.0] * NUM_MACHINES
    job_finish = [0.0] * NUM_JOBS
    ops_schedule = []  # list of dicts: {job, op_idx, machine, start, dur}
    for j_idx in range(NUM_JOBS):
        for op_idx in range(OPS_PER_JOB):
            dur = durations[j_idx][op_idx]
            m = required_machine_for_op(op_idx)
            earliest = max(machine_available[m], job_finish[j_idx])
            start = float(earliest)
            finish = start + dur
            # record
            ops_schedule.append({
                "job_index": j_idx,
                "job": JOBS[j_idx],
                "op": op_idx+1,
                "machine_index": m,
                "machine": MACHINES[m],
                "start": start,
                "duration": dur
            })
            # update
            machine_available[m] = finish
            job_finish[j_idx] = finish
    return ops_schedule

# ---------- Draw slide ----------
def create_pptx(ops_schedule):
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(LEFT_MARGIN), Inches(0.3), Inches(10), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "FJSP Gantt (3 machines × 5 jobs × 3 ops) — operations as separate shapes"
    p.font.size = Pt(18)

    # Draw machine labels and horizontal rows
    for m_idx, machine in enumerate(MACHINES):
        top = TOP_MARGIN + TITLE_HEIGHT + m_idx * (ROW_HEIGHT + ROW_GAP)
        # machine label
        lab = slide.shapes.add_textbox(Inches(LEFT_MARGIN - 0.9), Inches(top), Inches(0.9), Inches(ROW_HEIGHT))
        lab.text_frame.text = machine
        lab.text_frame.paragraphs[0].font.size = Pt(12)
        # row background (optional subtle line)
        row = slide.shapes.add_shape(
            1, Inches(top), Inches(SLIDE_WIDTH_IN - 2), Inches(ROW_HEIGHT)  # temporary params, will set properly below
        )
        # Instead of using shapes.add_shape for a row background (complex), we'll draw operations directly on time axis.

    # Draw time axis ticks (compute total horizon)
    last_end = max([op['start'] + op['duration'] for op in ops_schedule])
    horizon = last_end + 1
    # draw ticks
    for t in range(int(horizon)+1):
        x_pos = TIME_AXIS_OFFSET + t * TIME_SCALE
        # vertical line - use a thin rectangle to emulate a line
        slide.shapes.add_shape(
            1, Inches(x_pos), Inches(TOP_MARGIN + TITLE_HEIGHT - 0.1), Inches(0.02), Inches(NUM_MACHINES*(ROW_HEIGHT+ROW_GAP)+0.2)
        ).fill.background()  # invisible fallback; we will not rely on these lines in this script

    # Draw each operation as a rectangle shape positioned by start & duration
    for op in ops_schedule:
        job_idx = op['job_index']
        m_idx = op['machine_index']
        start = op['start']
        dur = op['duration']
        top = TOP_MARGIN + TITLE_HEIGHT + m_idx * (ROW_HEIGHT + ROW_GAP)
        left = TIME_AXIS_OFFSET + start * TIME_SCALE
        width = dur * TIME_SCALE
        height = ROW_HEIGHT * 0.8
        # Add rectangle
        shape = slide.shapes.add_shape(
            1, Inches(left), Inches(top + 0.08), Inches(width), Inches(height)
        )
        # Set fill color for the job
        r, g, b = JOB_COLORS[job_idx % len(JOB_COLORS)]
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(r, g, b)
        line = shape.line
        line.color.rgb = RGBColor(0, 0, 0)  # border
        # Add text label (job + op)
        text = f"{op['job']}-O{op['op']}\n{op['start']}→{op['start']+op['duration']}"
        tf = shape.text_frame
        tf.text = text
        for p in tf.paragraphs:
            p.font.size = Pt(10)

    # Legend (job colors)
    legend_left = TIME_AXIS_OFFSET + (horizon + 0.2) * TIME_SCALE
    lx = Inches(legend_left)
    ly = Inches(TOP_MARGIN + TITLE_HEIGHT)
    for j_idx, job in enumerate(JOBS):
        ly_j = TOP_MARGIN + TITLE_HEIGHT + j_idx * 0.28
        box = slide.shapes.add_shape(1, Inches(legend_left), Inches(ly_j), Inches(0.25), Inches(0.2))
        r, g, b = JOB_COLORS[j_idx % len(JOB_COLORS)]
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(r, g, b)
        box.line.color.rgb = RGBColor(0,0,0)
        tx = slide.shapes.add_textbox(Inches(legend_left+0.3), Inches(ly_j), Inches(1.8), Inches(0.22))
        tx.text_frame.text = job
        tx.text_frame.paragraphs[0].font.size = Pt(10)

    # Save
    fn = "fjsp_gantt.pptx"
    prs.save(fn)
    print(f"Saved {fn} with {len(ops_schedule)} operation shapes on one slide.")

if __name__ == "__main__":
    ops = schedule_jobs()
    create_pptx(ops)
